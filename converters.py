"""
iConvert conversion engine (no GUI dependencies).
Each converter takes absolute source + destination paths.
Platform/heavy libraries are imported lazily so this module always imports.
"""

import os
import tempfile
import shutil


# Target label -> (output extension, set of valid source extensions)
TARGETS = {
    "PDF":                ("pdf",  {"doc", "docx", "ppt", "pptx", "jpg", "jpeg", "png"}),
    "Word (.docx)":       ("docx", {"pdf"}),
    "PowerPoint (.pptx)": ("pptx", {"pdf"}),
    "PNG":                ("png",  {"jpg", "jpeg"}),
    "JPG":                ("jpg",  {"png"}),
}
TARGET_ORDER = ["PDF", "Word (.docx)", "PowerPoint (.pptx)", "PNG", "JPG"]
SUPPORTED_INPUTS = {"pdf", "doc", "docx", "ppt", "pptx", "jpg", "jpeg", "png"}


def unique_path(path):
    """If 'path' exists, append ' (1)', ' (2)', ... so we never overwrite."""
    if not os.path.exists(path):
        return path
    root, ext = os.path.splitext(path)
    i = 1
    while True:
        candidate = "{} ({}){}".format(root, i, ext)
        if not os.path.exists(candidate):
            return candidate
        i += 1


# --- image conversions (Pillow) --------------------------------------------
def convert_image(src, dst, target_ext):
    from PIL import Image
    img = Image.open(src)
    target_ext = target_ext.lower()
    if target_ext in ("jpg", "jpeg"):
        if img.mode in ("RGBA", "LA", "P"):           # JPEG has no transparency
            img = img.convert("RGBA")
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1])
            img = bg
        else:
            img = img.convert("RGB")
        img.save(dst, "JPEG", quality=95)
    elif target_ext == "png":
        img.save(dst, "PNG")
    else:
        raise ValueError("Unsupported image target: " + target_ext)


def image_to_pdf(src, dst):
    from PIL import Image
    Image.open(src).convert("RGB").save(dst, "PDF", resolution=150.0)


# --- PDF -> Word (pdf2docx) -------------------------------------------------
def pdf_to_word(src, dst):
    from pdf2docx import Converter
    cv = Converter(src)
    try:
        cv.convert(dst, start=0, end=None)
    finally:
        cv.close()


# --- PDF -> PowerPoint (PyMuPDF render + python-pptx) -----------------------
def pdf_to_ppt(src, dst):
    try:
        import fitz  # PyMuPDF
    except ImportError:
        import pymupdf as fitz
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(src)
    tmpdir = tempfile.mkdtemp(prefix="iconvert_")
    try:
        prs = Presentation()
        first = doc[0].rect                            # points -> EMU
        prs.slide_width = Emu(int(first.width * 12700))
        prs.slide_height = Emu(int(first.height * 12700))
        blank = prs.slide_layouts[6]
        zoom = 2.0

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            png = os.path.join(tmpdir, "page_{}.png".format(i))
            pix.save(png)
            slide = prs.slides.add_slide(blank)
            sw, sh = int(prs.slide_width), int(prs.slide_height)
            iw, ih = pix.width, pix.height
            if iw / ih > sw / sh:                       # aspect-fit, centered
                w = sw; h = int(sw * ih / iw); left = 0; top = int((sh - h) / 2)
            else:
                h = sh; w = int(sh * iw / ih); top = 0; left = int((sw - w) / 2)
            slide.shapes.add_picture(png, left, top, width=w, height=h)
        prs.save(dst)
    finally:
        doc.close()
        shutil.rmtree(tmpdir, ignore_errors=True)


# --- Word -> PDF (Microsoft Word via COM) ----------------------------------
def word_to_pdf(src, dst):
    import win32com.client
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    try:
        doc = word.Documents.Open(src, ReadOnly=True)
        doc.SaveAs(dst, FileFormat=17)                 # 17 = wdFormatPDF
        doc.Close(False)
    finally:
        word.Quit()


# --- PowerPoint -> PDF (Microsoft PowerPoint via COM) ----------------------
def ppt_to_pdf(src, dst):
    import win32com.client
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    try:
        deck = ppt.Presentations.Open(src, WithWindow=False)
        deck.SaveAs(dst, 32)                           # 32 = ppSaveAsPDF
        deck.Close()
    finally:
        ppt.Quit()


def route(src_ext, out_ext):
    """Return the conversion function for a (source, target) pair, or None."""
    src_ext = src_ext.lower()
    out_ext = out_ext.lower()
    if out_ext == "pdf":
        if src_ext in ("doc", "docx"):
            return word_to_pdf
        if src_ext in ("ppt", "pptx"):
            return ppt_to_pdf
        if src_ext in ("jpg", "jpeg", "png"):
            return image_to_pdf
    elif out_ext == "docx" and src_ext == "pdf":
        return pdf_to_word
    elif out_ext == "pptx" and src_ext == "pdf":
        return pdf_to_ppt
    elif out_ext in ("png", "jpg", "jpeg"):
        if src_ext in ("jpg", "jpeg", "png"):
            return lambda s, d: convert_image(s, d, out_ext)
    return None


# ---------------------------------------------------------------------------
# Version helpers (used by the in-app updater)
# ---------------------------------------------------------------------------
def parse_version(s):
    parts = []
    for p in str(s).strip().split("."):
        p = p.strip()
        try:
            parts.append(int(p))
        except ValueError:
            digits = "".join(ch for ch in p if ch.isdigit())
            parts.append(int(digits) if digits else 0)
    return tuple(parts) if parts else (0,)


def is_newer(remote, local):
    """True if remote version string is strictly newer than local."""
    return parse_version(remote) > parse_version(local)


# ===========================================================================
# v3 tools: batch, PDF merge/split/compress, PDF->images, Excel, password, OCR
# ===========================================================================

def gather_files(paths, exts):
    """Expand a mix of files/folders into a flat list of matching files."""
    exts = {e.lower().lstrip(".") for e in exts}
    found = []
    for p in paths:
        if os.path.isdir(p):
            for root, _dirs, files in os.walk(p):
                for fn in files:
                    if "." in fn and fn.rsplit(".", 1)[-1].lower() in exts:
                        found.append(os.path.join(root, fn))
        elif os.path.isfile(p):
            if "." in p and p.rsplit(".", 1)[-1].lower() in exts:
                found.append(p)
    return found


def pdf_merge(srcs, dst):
    import fitz
    out = fitz.open()
    try:
        for s in srcs:
            d = fitz.open(s)
            out.insert_pdf(d)
            d.close()
        out.save(dst)
    finally:
        out.close()
    return dst


def _parse_pages(spec, n):
    """'1-3,5,8-' -> sorted unique 0-based indices within [0,n). None if blank."""
    spec = (spec or "").strip()
    if not spec:
        return None
    out = []
    for part in spec.replace(" ", "").split(","):
        if not part:
            continue
        if "-" in part:
            a, _, b = part.partition("-")
            start = int(a) if a else 1
            end = int(b) if b else n
            for p in range(start, end + 1):
                if 1 <= p <= n:
                    out.append(p - 1)
        else:
            p = int(part)
            if 1 <= p <= n:
                out.append(p - 1)
    seen, res = set(), []
    for i in out:
        if i not in seen:
            seen.add(i)
            res.append(i)
    return res


def pdf_split(src, out_dir, base, pages_spec):
    import fitz
    doc = fitz.open(src)
    n = doc.page_count
    outputs = []
    try:
        idxs = _parse_pages(pages_spec, n)
        if idxs is None:
            for i in range(n):
                w = fitz.open()
                w.insert_pdf(doc, from_page=i, to_page=i)
                p = unique_path(os.path.join(out_dir, "%s_p%d.pdf" % (base, i + 1)))
                w.save(p)
                w.close()
                outputs.append(p)
        else:
            if not idxs:
                raise ValueError("no valid pages in '%s'" % pages_spec)
            w = fitz.open()
            for i in idxs:
                w.insert_pdf(doc, from_page=i, to_page=i)
            p = unique_path(os.path.join(out_dir, "%s_selected.pdf" % base))
            w.save(p)
            w.close()
            outputs.append(p)
    finally:
        doc.close()
    return outputs


def pdf_compress(src, dst):
    import fitz
    doc = fitz.open(src)
    try:
        doc.save(dst, garbage=4, deflate=True, clean=True)
    finally:
        doc.close()
    return dst


def pdf_to_images(src, out_dir, base, fmt, dpi=150, progress_cb=None):
    import fitz
    ext = "jpg" if str(fmt).lower() in ("jpg", "jpeg") else "png"
    doc = fitz.open(src)
    outputs = []
    mat = fitz.Matrix(dpi / 72.0, dpi / 72.0)
    try:
        total = doc.page_count
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            p = unique_path(os.path.join(out_dir, "%s_p%d.%s" % (base, i + 1, ext)))
            pix.save(p)
            outputs.append(p)
            if progress_cb:
                progress_cb(i + 1, total)
    finally:
        doc.close()
    return outputs


def excel_to_pdf(src, dst):
    import win32com.client
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    try:
        wb = excel.Workbooks.Open(src, ReadOnly=True)
        wb.ExportAsFixedFormat(0, dst)  # 0 = xlTypePDF
        wb.Close(False)
    finally:
        excel.Quit()
    return dst


def excel_to_csv(src, out_dir, base):
    import openpyxl
    import csv
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    outputs = []
    names = wb.sheetnames
    try:
        for name in names:
            ws = wb[name]
            if len(names) == 1:
                p = unique_path(os.path.join(out_dir, base + ".csv"))
            else:
                safe = "".join(c if (c.isalnum() or c in " -_") else "_" for c in name)
                p = unique_path(os.path.join(out_dir, "%s_%s.csv" % (base, safe)))
            with open(p, "w", newline="", encoding="utf-8-sig") as f:
                w = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    w.writerow(["" if c is None else c for c in row])
            outputs.append(p)
    finally:
        wb.close()
    return outputs


def pdf_encrypt(src, dst, password):
    import fitz
    if not password:
        raise ValueError("password is empty")
    doc = fitz.open(src)
    try:
        perm = int(fitz.PDF_PERM_ACCESSIBILITY | fitz.PDF_PERM_PRINT |
                   fitz.PDF_PERM_COPY | fitz.PDF_PERM_ANNOTATE)
        doc.save(dst, encryption=fitz.PDF_ENCRYPT_AES_256,
                 owner_pw=password, user_pw=password, permissions=perm)
    finally:
        doc.close()
    return dst


def pdf_decrypt(src, dst, password):
    import fitz
    doc = fitz.open(src)
    try:
        if doc.needs_pass:
            if not doc.authenticate(password or ""):
                raise ValueError("wrong password")
        doc.save(dst, encryption=fitz.PDF_ENCRYPT_NONE)
    finally:
        doc.close()
    return dst


def _find_tesseract():
    """Locate tesseract.exe even when it is not on PATH (the Windows installer
    usually does NOT add it to PATH)."""
    import shutil
    env = os.environ.get("TESSERACT_CMD")
    if env and os.path.isfile(env):
        return env
    onpath = shutil.which("tesseract")
    if onpath:
        return onpath
    la = os.environ.get("LOCALAPPDATA", "")
    roots = [
        os.environ.get("ProgramFiles", "C:/Program Files"),
        os.environ.get("ProgramFiles(x86)", "C:/Program Files (x86)"),
        os.path.join(la, "Programs") if la else "",
        la,
    ]
    for r in roots:
        if not r:
            continue
        c = os.path.join(r, "Tesseract-OCR", "tesseract.exe")
        try:
            if os.path.isfile(c):
                return c
        except Exception:
            pass
    return None


def _configure_tesseract():
    try:
        import pytesseract
    except Exception:
        return None
    path = _find_tesseract()
    if path:
        try:
            pytesseract.pytesseract.tesseract_cmd = path
        except Exception:
            pass
    return path


def tesseract_ok():
    try:
        import pytesseract
        _configure_tesseract()
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def pdf_ocr(src, dst, dpi=200, progress_cb=None):
    import io
    import fitz
    import pytesseract
    from PIL import Image
    _configure_tesseract()
    doc = fitz.open(src)
    out = fitz.open()
    mat = fitz.Matrix(dpi / 72.0, dpi / 72.0)
    try:
        total = doc.page_count
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            data = pytesseract.image_to_pdf_or_hocr(img, extension="pdf")
            sub = fitz.open("pdf", data)
            out.insert_pdf(sub)
            sub.close()
            if progress_cb:
                progress_cb(i + 1, total)
        out.save(dst)
    finally:
        out.close()
        doc.close()
    return dst


def pdf_ocr_to_word(src, dst, dpi=200, progress_cb=None):
    import io
    import fitz
    import pytesseract
    from PIL import Image
    _configure_tesseract()
    from docx import Document
    doc = fitz.open(src)
    out = Document()
    mat = fitz.Matrix(dpi / 72.0, dpi / 72.0)
    try:
        total = doc.page_count
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            text = pytesseract.image_to_string(img)
            if i > 0:
                out.add_page_break()
            for line in text.splitlines():
                out.add_paragraph(line)
            if progress_cb:
                progress_cb(i + 1, total)
    finally:
        doc.close()
    out.save(dst)
    return dst


# ===========================================================================
# v3.2 tools: rotate / delete pages / page numbers / watermark / extract
# images / images->PDF / image resize / OCR language list
# ===========================================================================

def pdf_rotate(src, dst, degrees):
    import fitz
    d = int(degrees) % 360
    doc = fitz.open(src)
    try:
        for page in doc:
            page.set_rotation((page.rotation + d) % 360)
        doc.save(dst)
    finally:
        doc.close()
    return dst


def pdf_delete_pages(src, out_dir, base, pages_spec):
    import fitz
    doc = fitz.open(src)
    try:
        n = doc.page_count
        idxs = _parse_pages(pages_spec, n)
        if not idxs:
            raise ValueError("say which pages to remove, e.g. 2,5")
        drop = set(idxs)
        keep = [i for i in range(n) if i not in drop]
        if not keep:
            raise ValueError("that would remove every page")
        doc.select(keep)
        p = unique_path(os.path.join(out_dir, base + "_edited.pdf"))
        doc.save(p)
    finally:
        doc.close()
    return [p]


def pdf_add_page_numbers(src, dst):
    import fitz
    doc = fitz.open(src)
    try:
        n = doc.page_count
        for i, page in enumerate(doc, 1):
            r = page.rect
            page.insert_text((r.width / 2.0 - 18, r.height - 26),
                             "%d / %d" % (i, n), fontsize=10, color=(0.35, 0.35, 0.35))
        doc.save(dst)
    finally:
        doc.close()
    return dst


def pdf_watermark(src, dst, text):
    import fitz
    if not text:
        raise ValueError("enter watermark text")
    doc = fitz.open(src)
    try:
        for page in doc:
            r = page.rect
            page.insert_text((r.width * 0.12, r.height * 0.52), text,
                             fontsize=44, color=(0.82, 0.82, 0.82))
        doc.save(dst)
    finally:
        doc.close()
    return dst


def pdf_extract_images(src, out_dir, base):
    import fitz
    doc = fitz.open(src)
    outs = []
    try:
        seen = set()
        count = 0
        for pno in range(doc.page_count):
            for img in doc.get_page_images(pno):
                xref = img[0]
                if xref in seen:
                    continue
                seen.add(xref)
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha >= 4:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                count += 1
                p = unique_path(os.path.join(out_dir, "%s_img%d.png" % (base, count)))
                pix.save(p)
                pix = None
                outs.append(p)
    finally:
        doc.close()
    if not outs:
        raise ValueError("no embedded images found in this PDF")
    return outs


def images_to_pdf(srcs, dst):
    from PIL import Image
    imgs = []
    for sp in srcs:
        im = Image.open(sp)
        imgs.append(im.convert("RGB"))
    if not imgs:
        raise ValueError("no images")
    imgs[0].save(dst, "PDF", save_all=True, append_images=imgs[1:], resolution=150.0)
    return dst


def image_resize(src, dst, max_side):
    from PIL import Image
    max_side = int(max_side)
    if max_side < 1:
        raise ValueError("max size must be a positive number")
    im = Image.open(src)
    w, h = im.size
    scale = min(1.0, float(max_side) / float(max(w, h)))
    if scale < 1.0:
        im = im.resize((max(1, int(w * scale)), max(1, int(h * scale))), Image.LANCZOS)
    ext = os.path.splitext(dst)[1].lower()
    if ext in (".jpg", ".jpeg"):
        if im.mode in ("RGBA", "LA", "P"):
            im = im.convert("RGBA")
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[-1])
            im = bg
        else:
            im = im.convert("RGB")
        im.save(dst, "JPEG", quality=85)
    else:
        im.save(dst)
    return dst


def list_ocr_langs():
    try:
        import pytesseract
        _configure_tesseract()
        return sorted(pytesseract.get_languages(config=""))
    except Exception:
        return []


# ===========================================================================
# Batch runner used both in-thread and in a separate process (never-freeze)
# ===========================================================================

def convert_one(tool, src, out_dir, opts, progress_cb=None):
    op = tool["op"]
    base = os.path.splitext(os.path.basename(src))[0]

    def d(ext, suffix=""):
        return unique_path(os.path.join(out_dir, base + suffix + "." + ext))

    if op == "img":
        o = d(tool["out"]); convert_image(src, o, tool["out"]); return [o]
    if op == "img2pdf":
        o = d("pdf"); image_to_pdf(src, o); return [o]
    if op == "pdf2word":
        o = d("docx"); pdf_to_word(src, o); return [o]
    if op == "pdf2ppt":
        o = d("pptx"); pdf_to_ppt(src, o); return [o]
    if op == "word2pdf":
        o = d("pdf"); word_to_pdf(src, o); return [o]
    if op == "ppt2pdf":
        o = d("pdf"); ppt_to_pdf(src, o); return [o]
    if op == "excel2pdf":
        o = d("pdf"); excel_to_pdf(src, o); return [o]
    if op == "excel2csv":
        return excel_to_csv(src, out_dir, base)
    if op == "pdf2img":
        return pdf_to_images(src, out_dir, base, tool["fmt"], progress_cb=progress_cb)
    if op == "compress":
        o = d("pdf", "_compressed"); pdf_compress(src, o); return [o]
    if op == "split":
        return pdf_split(src, out_dir, base, opts.get("pages", ""))
    if op == "encrypt":
        o = d("pdf", "_protected"); pdf_encrypt(src, o, opts.get("password", "")); return [o]
    if op == "decrypt":
        o = d("pdf", "_unlocked"); pdf_decrypt(src, o, opts.get("password", "")); return [o]
    if op == "ocr_pdf":
        o = d("pdf", "_ocr"); pdf_ocr(src, o, progress_cb=progress_cb); return [o]
    if op == "ocr_word":
        o = d("docx", "_ocr"); pdf_ocr_to_word(src, o, progress_cb=progress_cb); return [o]
    if op == "rotate":
        o = d("pdf", "_rotated"); pdf_rotate(src, o, opts.get("degrees", 90)); return [o]
    if op == "delpages":
        return pdf_delete_pages(src, out_dir, base, opts.get("pages", ""))
    if op == "pagenum":
        o = d("pdf", "_numbered"); pdf_add_page_numbers(src, o); return [o]
    if op == "watermark":
        o = d("pdf", "_watermark"); pdf_watermark(src, o, opts.get("text", "")); return [o]
    if op == "extractimg":
        return pdf_extract_images(src, out_dir, base)
    if op == "resize":
        ext = os.path.splitext(src)[1].lstrip(".").lower() or "png"
        o = d(ext, "_small"); image_resize(src, o, opts.get("number", 1600)); return [o]
    raise ValueError("unknown op " + op)


def process_batch(tool, files, out_dir, opts, q):
    """Run a whole batch, posting progress messages onto q. Safe to run in a
    separate process (no GUI imports)."""
    com = False
    try:
        import pythoncom
        pythoncom.CoInitialize()
        com = True
    except Exception:
        pass
    ok = 0
    total = 1
    try:
        if tool.get("kind") == "combine":
            total = 1
            q.put(("start", (1, 1, "Combining %d files" % len(files))))
            try:
                folder = out_dir or os.path.dirname(files[0])
                if tool["op"] == "images_pdf":
                    dst = unique_path(os.path.join(folder, "combined.pdf"))
                    images_to_pdf(files, dst)
                else:
                    dst = unique_path(os.path.join(folder, "merged.pdf"))
                    pdf_merge(files, dst)
                ok = 1
                q.put(("lastdir", folder))
                q.put(("result", (True, "%d files" % len(files), os.path.basename(dst))))
            except Exception as e:
                q.put(("result", (False, tool.get("title", "combine"), str(e)[:80])))
            q.put(("doneone", (1, 1)))
        else:
            total = len(files)
            for i, src in enumerate(files, start=1):
                q.put(("start", (i, total, os.path.basename(src))))

                def cb(done, pages):
                    q.put(("sub", (done, pages)))

                try:
                    folder = out_dir or os.path.dirname(src)
                    outs = convert_one(tool, src, folder, opts, progress_cb=cb)
                    ok += 1
                    q.put(("lastdir", folder))
                    label = (os.path.basename(outs[0]) if len(outs) == 1
                             else "%d files" % len(outs))
                    q.put(("result", (True, os.path.basename(src), label)))
                except Exception as e:
                    q.put(("result", (False, os.path.basename(src), str(e)[:80])))
                q.put(("doneone", (i, total)))
    finally:
        if com:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
    q.put(("alldone", (ok, total)))
